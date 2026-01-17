#!/usr/bin/env python3
"""
Content Generation Tools for Jotty
Ported from JustJot.ai adapters/sinks/

Provides PDF, HTML, Markdown, DOCX, and PPTX generation
"""

from pathlib import Path
from typing import Optional, Dict, Any
import subprocess
import re
from datetime import datetime

from .document import Document, Section, SectionType

# Optional dependencies
try:
    from docx import Document as DocxDocument
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class ContentGenerators:
    """
    Content generation tools for Jotty agents

    Provides methods to create PDFs, HTML, and Markdown files from documents
    """

    def __init__(self):
        self.output_dir = Path("./outputs/research")
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def _sanitize_filename(self, filename: str) -> str:
        """Sanitize filename to be safe for filesystem"""
        # Remove/replace unsafe characters
        filename = re.sub(r'[<>:"/\\|?*]', '', filename)
        # Replace spaces with underscores
        filename = filename.replace(' ', '_')
        # Remove leading/trailing dots and spaces
        filename = filename.strip('. ')
        # Limit length
        filename = filename[:200]
        return filename if filename else 'document'

    # =========================================================================
    # PDF Generation (via Pandoc + XeLaTeX)
    # =========================================================================

    def generate_pdf(
        self,
        document: Document,
        output_path: Optional[Path] = None,
        format: str = 'a4'
    ) -> Path:
        """
        Generate PDF from document using pandoc + XeLaTeX

        Args:
            document: Document to convert
            output_path: Output directory (default: ./outputs/research)
            format: PDF format (a4, a5, letter) - default: a4

        Returns:
            Path to generated PDF

        Raises:
            RuntimeError: If conversion fails
        """
        # Format mapping
        FORMAT_MAP = {
            'a4': 'a4',
            'a5': 'a5',
            'a6': 'a6',
            'letter': 'letter',
        }

        pdf_format = format.lower()
        if pdf_format not in FORMAT_MAP:
            raise ValueError(f"Unknown PDF format: {pdf_format}. "
                           f"Must be one of: {list(FORMAT_MAP.keys())}")

        pandoc_page_size = FORMAT_MAP[pdf_format]

        # Determine output path
        if output_path is None:
            output_path = self.output_dir

        output_path = Path(output_path)
        output_path.mkdir(parents=True, exist_ok=True)

        # Generate filename
        filename = self._sanitize_filename(document.title)
        pdf_path = output_path / f"{filename}_{pdf_format}.pdf"

        # Get markdown content
        markdown_content = document.full_content
        if not markdown_content or not markdown_content.strip():
            raise ValueError("Document has no content to convert")

        # Save temporary markdown file
        temp_md = output_path / f"{filename}_temp.md"
        try:
            with open(temp_md, 'w', encoding='utf-8') as f:
                f.write(f"# {document.title}\n\n")
                if document.author:
                    f.write(f"**Author:** {document.author}\n\n")
                f.write(f"**Generated:** {document.created.strftime('%Y-%m-%d')}\n\n")
                f.write("---\n\n")
                f.write(markdown_content)

            print(f"      Converting to PDF ({pdf_format.upper()})...")

            # Use pandoc to convert to PDF
            cmd = [
                'pandoc',
                str(temp_md),
                '-f', 'markdown',
                '-t', 'pdf',
                '--pdf-engine=xelatex',
                f'--variable=papersize:{pandoc_page_size}',
                '-o', str(pdf_path)
            ]

            result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)

            if result.returncode != 0:
                error_output = result.stderr if result.stderr else result.stdout
                raise RuntimeError(f"Pandoc failed with exit code {result.returncode}: {error_output}")

        except subprocess.TimeoutExpired:
            raise RuntimeError("Pandoc conversion timed out (>60 seconds)")
        except Exception as e:
            raise RuntimeError(f"PDF conversion failed: {e}")
        finally:
            # Always try to clean up temp file
            try:
                if temp_md.exists():
                    temp_md.unlink()
            except Exception as cleanup_error:
                print(f"      Warning: Failed to clean up temp file '{temp_md}': {cleanup_error}")

        # Verify output file was created
        if not pdf_path.exists():
            raise RuntimeError(
                f"PDF conversion completed but output file was not created at '{pdf_path}'. "
                f"This may indicate a silent failure in pandoc or the LaTeX engine."
            )

        # Verify output file has content
        file_size = pdf_path.stat().st_size
        if file_size == 0:
            pdf_path.unlink()
            raise RuntimeError(f"Generated PDF file is empty (0 bytes) - conversion failed")

        print(f"      ✅ PDF created: {pdf_path.name} ({file_size} bytes)")
        return pdf_path

    # =========================================================================
    # HTML Generation (via Pandoc)
    # =========================================================================

    def generate_html(
        self,
        document: Document,
        output_path: Optional[Path] = None,
        standalone: bool = True,
        include_toc: bool = True
    ) -> Path:
        """
        Generate HTML from document using pandoc

        Args:
            document: Document to convert
            output_path: Output directory (default: ./outputs/research)
            standalone: Create standalone HTML with CSS (default: True)
            include_toc: Include table of contents (default: True)

        Returns:
            Path to generated HTML
        """
        # Determine output path
        if output_path is None:
            output_path = self.output_dir

        output_path = Path(output_path)
        output_path.mkdir(parents=True, exist_ok=True)

        # Generate filename
        filename = self._sanitize_filename(document.title)
        html_path = output_path / f"{filename}.html"

        # Get markdown content
        markdown_content = document.full_content

        # Save temporary markdown file
        temp_md = output_path / f"{filename}_temp.md"
        try:
            with open(temp_md, 'w', encoding='utf-8') as f:
                f.write(f"# {document.title}\n\n")
                if document.author:
                    f.write(f"**Author:** {document.author}\n\n")
                f.write(f"**Generated:** {document.created.strftime('%Y-%m-%d')}\n\n")
                f.write("---\n\n")
                f.write(markdown_content)

            print(f"      Converting to HTML...")

            # Use pandoc to convert markdown to HTML
            cmd = [
                'pandoc',
                str(temp_md),
                '-o', str(html_path),
                '--mathml',  # Use MathML for math
            ]

            # Add standalone HTML with CSS if requested
            if standalone:
                cmd.append('--standalone')
                cmd.append('--self-contained')

            # Add table of contents if requested
            if include_toc:
                cmd.extend(['--toc', '--toc-depth=3'])

            # Add metadata
            cmd.extend([
                '--metadata', f'title={document.title}',
            ])

            if document.author:
                cmd.extend(['--metadata', f'author={document.author}'])

            # Run pandoc
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=60
            )

            if result.returncode != 0:
                raise RuntimeError(f"Pandoc failed: {result.stderr}")

            # Clean up temp file
            temp_md.unlink()

            print(f"      ✅ HTML: {html_path.name}")

            return html_path

        except FileNotFoundError:
            raise RuntimeError("pandoc not found. Please install pandoc: https://pandoc.org/installing.html")
        except subprocess.TimeoutExpired:
            raise RuntimeError("HTML conversion timed out")
        except Exception as e:
            # Clean up temp file on error
            if temp_md.exists():
                temp_md.unlink()
            raise RuntimeError(f"HTML conversion failed: {e}")

    # =========================================================================
    # Markdown Export
    # =========================================================================

    def export_markdown(
        self,
        document: Document,
        output_path: Optional[Path] = None,
        include_metadata: bool = True
    ) -> Path:
        """
        Export document to Markdown file

        Args:
            document: Document to export
            output_path: Output directory (default: ./outputs/research)
            include_metadata: Include YAML frontmatter (default: True)

        Returns:
            Path to generated .md file
        """
        # Determine output path
        if output_path is None:
            output_path = self.output_dir

        output_path = Path(output_path)
        output_path.mkdir(parents=True, exist_ok=True)

        # Generate filename
        filename = self._sanitize_filename(document.title)
        date_prefix = document.created.strftime('%Y-%m-%d')
        md_path = output_path / f"{date_prefix}-{filename}.md"

        print(f"      Exporting to Markdown...")

        # Build markdown content
        md_content = []

        if include_metadata:
            # Add frontmatter metadata
            md_content.append("---")
            yaml_title = document.title.replace('"', '\\"')
            md_content.append(f'title: "{yaml_title}"')
            if document.author:
                md_content.append(f"author: {document.author}")
            if document.topic:
                md_content.append(f"topic: {document.topic}")
            md_content.append(f"date: {document.created.strftime('%Y-%m-%d')}")
            md_content.append(f"source: {document.source_type}")
            md_content.append("---")
            md_content.append("")

        # Add title
        md_content.append(f"# {document.title}")
        md_content.append("")

        # Add content
        md_content.append(document.full_content)

        # Write file
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(md_content))

        print(f"      ✅ MD: {md_path.name}")

        return md_path

    # =========================================================================
    # DOCX Generation (via python-docx)
    # =========================================================================

    def generate_docx(
        self,
        document: Document,
        output_path: Optional[Path] = None
    ) -> Path:
        """
        Generate Word document from Document

        Args:
            document: Document to convert
            output_path: Output directory (default: ./outputs/research)

        Returns:
            Path to generated .docx file

        Raises:
            RuntimeError: If python-docx not installed
        """
        if not HAS_DOCX:
            raise RuntimeError(
                "python-docx not installed. Install with: pip install python-docx"
            )

        # Determine output path
        if output_path is None:
            output_path = self.output_dir

        output_path = Path(output_path)
        output_path.mkdir(parents=True, exist_ok=True)

        # Generate filename
        filename = self._sanitize_filename(document.title)
        docx_path = output_path / f"{filename}.docx"

        print(f"      Generating DOCX...")

        # Create document
        doc = DocxDocument()

        # Add title
        title = doc.add_heading(document.title, level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Add metadata
        if document.author:
            author_para = doc.add_paragraph(f"Author: {document.author}")
            author_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            author_para.runs[0].italic = True

        date_para = doc.add_paragraph(f"Date: {document.created.strftime('%Y-%m-%d')}")
        date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        date_para.runs[0].italic = True

        doc.add_paragraph()  # Spacing

        # Add sections
        for section in document.sections:
            if section.title:
                doc.add_heading(section.title, level=1)

            if section.type == SectionType.TEXT:
                # Add paragraphs
                for para in section.content.split('\n\n'):
                    if para.strip():
                        doc.add_paragraph(para.strip())

            elif section.type == SectionType.CODE:
                # Add code block
                code_para = doc.add_paragraph(section.content)
                code_para.style = 'No Spacing'
                for run in code_para.runs:
                    run.font.name = 'Courier New'
                    run.font.size = Pt(10)

            elif section.type == SectionType.MATH:
                # Math as preformatted text (DOCX doesn't support LaTeX natively)
                math_para = doc.add_paragraph(section.content)
                math_para.style = 'Intense Quote'

            elif section.type == SectionType.MERMAID:
                # Mermaid diagrams as code blocks
                doc.add_paragraph("Diagram:", style='Intense Quote')
                diagram_para = doc.add_paragraph(section.content)
                diagram_para.style = 'No Spacing'

            else:
                # Default: plain text
                doc.add_paragraph(section.content)

        # If no sections, use flat content
        if not document.sections and document.content:
            for para in document.content.split('\n\n'):
                if para.strip():
                    doc.add_paragraph(para.strip())

        # Save
        doc.save(docx_path)

        print(f"      ✅ DOCX: {docx_path.name}")

        return docx_path

    # =========================================================================
    # PPTX Generation (via python-pptx)
    # =========================================================================

    def generate_pptx(
        self,
        document: Document,
        output_path: Optional[Path] = None
    ) -> Path:
        """
        Generate PowerPoint presentation from Document

        Args:
            document: Document to convert
            output_path: Output directory (default: ./outputs/research)

        Returns:
            Path to generated .pptx file

        Raises:
            RuntimeError: If python-pptx not installed
        """
        if not HAS_PPTX:
            raise RuntimeError(
                "python-pptx not installed. Install with: pip install python-pptx"
            )

        # Determine output path
        if output_path is None:
            output_path = self.output_dir

        output_path = Path(output_path)
        output_path.mkdir(parents=True, exist_ok=True)

        # Generate filename
        filename = self._sanitize_filename(document.title)
        pptx_path = output_path / f"{filename}.pptx"

        print(f"      Generating PPTX...")

        # Create presentation
        prs = Presentation()

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = document.title
        subtitle.text = f"{document.author or 'Jotty'}\n{document.created.strftime('%Y-%m-%d')}"

        # Content slides
        for section in document.sections:
            # Use title and content layout
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]

            title_shape.text = section.title or "Section"

            # Add content based on type
            if section.type == SectionType.TEXT:
                text_frame = body_shape.text_frame
                text_frame.text = section.content

            elif section.type == SectionType.CODE:
                text_frame = body_shape.text_frame
                text_frame.text = section.content
                # Make it monospace-like
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Courier New'
                        run.font.size = PptxPt(10)

            elif section.type == SectionType.MERMAID:
                text_frame = body_shape.text_frame
                text_frame.text = f"Diagram:\n\n{section.content}"

            else:
                text_frame = body_shape.text_frame
                text_frame.text = section.content

        # Save
        prs.save(pptx_path)

        print(f"      ✅ PPTX: {pptx_path.name}")

        return pptx_path
