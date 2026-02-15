"""
PowerPoint Editor Skill

Editing toolkit for existing PowerPoint (.pptx) files using python-pptx.
Also includes PptxGenJS-based presentation creation via Node.js scripts.
"""

import json
import logging
import os
import shutil
import subprocess
import uuid
import xml.etree.ElementTree as ET
import zipfile
from pathlib import Path
from typing import Any, Dict, List, Optional

from Jotty.core.infrastructure.utils.env_loader import load_jotty_env
from Jotty.core.infrastructure.utils.skill_status import SkillStatus
from Jotty.core.infrastructure.utils.tool_helpers import tool_error, tool_response, tool_wrapper

load_jotty_env()

# Status emitter for progress updates
status = SkillStatus("pptx-editor")

logger = logging.getLogger(__name__)

# Storage for generated presentations
PPTX_STORAGE_PATH = os.path.expanduser("~/jotty/presentations")


class PPTXEditorError(Exception):
    """Custom exception for PPTX editor operations."""

    pass


class PPTXEditor:
    """Class for handling PowerPoint editing operations."""

    @staticmethod
    def _validate_file_path(file_path: str) -> None:
        """Validate that file exists and is a PPTX file."""
        if not file_path:
            raise PPTXEditorError("file_path parameter is required")
        if not os.path.exists(file_path):
            raise PPTXEditorError(f"File not found: {file_path}")
        if not file_path.lower().endswith(".pptx"):
            raise PPTXEditorError("File must be a .pptx file")

    @staticmethod
    def _get_presentation(file_path: str):
        """Load and return a Presentation object."""
        try:
            from pptx import Presentation
        except ImportError:
            raise PPTXEditorError(
                "python-pptx not installed. Install with: pip install python-pptx"
            )
        return Presentation(file_path)

    @staticmethod
    def _extract_shape_text(shape) -> str:
        """Extract text from a shape."""
        if hasattr(shape, "text"):
            return shape.text
        if hasattr(shape, "text_frame"):
            return "\n".join([p.text for p in shape.text_frame.paragraphs])
        return ""

    @staticmethod
    def _get_slide_content(slide, include_notes: bool = False) -> Dict[str, Any]:
        """Extract content from a single slide."""
        content = {"title": "", "content": [], "shapes_count": len(slide.shapes)}

        for shape in slide.shapes:
            text = PPTXEditor._extract_shape_text(shape)
            if text:
                if shape.has_text_frame and shape == slide.shapes.title:
                    content["title"] = text
                else:
                    content["content"].append(text)

        if include_notes and slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else ""
            content["notes"] = notes_text

        return content


@tool_wrapper()
def read_pptx_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Read presentation content from a PPTX file.

    Args:
        params: Dictionary containing:
            - file_path (str, required): Path to the PPTX file
            - include_notes (bool, optional): Include slide notes (default: False)

    Returns:
        Dictionary with:
            - success (bool): Whether operation succeeded
            - slides (list): List of slide content dicts
            - slide_count (int): Total number of slides
            - error (str, optional): Error message if failed
    """
    status.set_callback(params.pop("_status_callback", None))

    file_path = params.get("file_path")
    include_notes = params.get("include_notes", False)

    try:
        PPTXEditor._validate_file_path(file_path)
        prs = PPTXEditor._get_presentation(file_path)

        slides = []
        for idx, slide in enumerate(prs.slides):
            slide_content = PPTXEditor._get_slide_content(slide, include_notes)
            slide_content["index"] = idx
            slides.append(slide_content)

        logger.info(f"Read {len(slides)} slides from: {file_path}")

        return {
            "success": True,
            "slides": slides,
            "slide_count": len(slides),
            "file_path": file_path,
        }

    except PPTXEditorError as e:
        return {"success": False, "error": str(e)}
    except Exception as e:
        logger.error(f"Failed to read PPTX: {e}", exc_info=True)
        return {"success": False, "error": f"Failed to read presentation: {str(e)}"}


@tool_wrapper()
def get_slide_layouts_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get available slide layouts from a PPTX file.

    Args:
        params: Dictionary containing:
            - file_path (str, required): Path to the PPTX file

    Returns:
        Dictionary with:
            - success (bool): Whether operation succeeded
            - layouts (list): List of layout dicts with index and name
            - error (str, optional): Error message if failed
    """
    status.set_callback(params.pop("_status_callback", None))

    file_path = params.get("file_path")

    try:
        PPTXEditor._validate_file_path(file_path)
        prs = PPTXEditor._get_presentation(file_path)

        layouts = []
        for idx, layout in enumerate(prs.slide_layouts):
            layouts.append({"index": idx, "name": layout.name})

        logger.info(f"Found {len(layouts)} layouts in: {file_path}")

        return {"success": True, "layouts": layouts, "layout_count": len(layouts)}

    except PPTXEditorError as e:
        return {"success": False, "error": str(e)}
    except Exception as e:
        logger.error(f"Failed to get layouts: {e}", exc_info=True)
        return {"success": False, "error": f"Failed to get layouts: {str(e)}"}


@tool_wrapper()
def add_slide_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Add a new slide to a presentation.

    Args:
        params: Dictionary containing:
            - file_path (str, required): Path to the PPTX file
            - layout_index (int, optional): Index of slide layout to use (default: 1)
            - title (str, optional): Slide title
            - content (str or list, optional): Slide content (text or list of bullets)
            - position (int, optional): Position to insert slide (default: end)

    Returns:
        Dictionary with:
            - success (bool): Whether operation succeeded
            - slide_index (int): Index of the new slide
            - slide_count (int): Total slides after addition
            - error (str, optional): Error message if failed
    """
    status.set_callback(params.pop("_status_callback", None))

    file_path = params.get("file_path")
    layout_index = params.get("layout_index", 1)
    title = params.get("title", "")
    content = params.get("content", "")
    position = params.get("position")

    try:
        from pptx.util import Inches, Pt

        PPTXEditor._validate_file_path(file_path)
        prs = PPTXEditor._get_presentation(file_path)

        if layout_index >= len(prs.slide_layouts):
            layout_index = 1 if len(prs.slide_layouts) > 1 else 0

        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)

        # Set title if shape exists
        if slide.shapes.title and title:
            slide.shapes.title.text = title

        # Add content
        if content:
            # Find content placeholder or create text box
            content_placeholder = None
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    content_placeholder = shape
                    break

            if content_placeholder:
                tf = content_placeholder.text_frame
                if isinstance(content, list):
                    for i, bullet in enumerate(content):
                        if i == 0:
                            tf.paragraphs[0].text = bullet
                        else:
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.level = 0
                else:
                    tf.text = content
            else:
                # Create text box
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                height = Inches(5)
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                tf.word_wrap = True

                if isinstance(content, list):
                    for i, bullet in enumerate(content):
                        if i == 0:
                            tf.paragraphs[0].text = f"- {bullet}"
                        else:
                            p = tf.add_paragraph()
                            p.text = f"- {bullet}"
                else:
                    tf.text = content

        # Move slide to position if specified
        if position is not None and position < len(prs.slides) - 1:
            # Get the slide element
            slide_elem = slide._element
            slides_elem = prs.slides._sldIdLst
            # Move to position
            slides_elem.remove(slides_elem[-1])
            slides_elem.insert(position, slide._element.getparent())

        prs.save(file_path)
        new_index = len(prs.slides) - 1 if position is None else position

        logger.info(f"Added slide at index {new_index} to: {file_path}")

        return {
            "success": True,
            "slide_index": new_index,
            "slide_count": len(prs.slides),
            "file_path": file_path,
        }

    except PPTXEditorError as e:
        return {"success": False, "error": str(e)}
    except Exception as e:
        logger.error(f"Failed to add slide: {e}", exc_info=True)
        return {"success": False, "error": f"Failed to add slide: {str(e)}"}


@tool_wrapper()
def update_slide_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Update content of an existing slide.

    Args:
        params: Dictionary containing:
            - file_path (str, required): Path to the PPTX file
            - slide_index (int, required): Index of slide to update (0-based)
            - title (str, optional): New slide title
            - content (str or list, optional): New slide content

    Returns:
        Dictionary with:
            - success (bool): Whether operation succeeded
            - slide_index (int): Index of updated slide
            - error (str, optional): Error message if failed
    """
    status.set_callback(params.pop("_status_callback", None))

    file_path = params.get("file_path")
    slide_index = params.get("slide_index")
    title = params.get("title")
    content = params.get("content")

    try:
        from pptx.util import Inches

        PPTXEditor._validate_file_path(file_path)

        if slide_index is None:
            return {"success": False, "error": "slide_index parameter is required"}

        prs = PPTXEditor._get_presentation(file_path)

        if slide_index < 0 or slide_index >= len(prs.slides):
            return {
                "success": False,
                "error": f"Invalid slide_index: {slide_index}. Presentation has {len(prs.slides)} slides.",
            }

        slide = prs.slides[slide_index]

        # Update title
        if title is not None and slide.shapes.title:
            slide.shapes.title.text = title

        # Update content
        if content is not None:
            content_placeholder = None
            for shape in slide.shapes:
                if shape.has_text_frame and shape != slide.shapes.title:
                    content_placeholder = shape
                    break

            if content_placeholder:
                tf = content_placeholder.text_frame
                # Clear existing paragraphs
                for para in tf.paragraphs:
                    para.clear()

                if isinstance(content, list):
                    for i, bullet in enumerate(content):
                        if i == 0:
                            tf.paragraphs[0].text = bullet
                        else:
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.level = 0
                else:
                    tf.paragraphs[0].text = content

        prs.save(file_path)

        logger.info(f"Updated slide {slide_index} in: {file_path}")

        return {"success": True, "slide_index": slide_index, "file_path": file_path}

    except PPTXEditorError as e:
        return {"success": False, "error": str(e)}
    except Exception as e:
        logger.error(f"Failed to update slide: {e}", exc_info=True)
        return {"success": False, "error": f"Failed to update slide: {str(e)}"}


@tool_wrapper()
def delete_slide_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Delete a slide from a presentation.

    Args:
        params: Dictionary containing:
            - file_path (str, required): Path to the PPTX file
            - slide_index (int, required): Index of slide to delete (0-based)

    Returns:
        Dictionary with:
            - success (bool): Whether operation succeeded
            - deleted_index (int): Index of deleted slide
            - slide_count (int): Total slides after deletion
            - error (str, optional): Error message if failed
    """
    status.set_callback(params.pop("_status_callback", None))

    file_path = params.get("file_path")
    slide_index = params.get("slide_index")

    try:
        PPTXEditor._validate_file_path(file_path)

        if slide_index is None:
            return {"success": False, "error": "slide_index parameter is required"}

        prs = PPTXEditor._get_presentation(file_path)

        if slide_index < 0 or slide_index >= len(prs.slides):
            return {
                "success": False,
                "error": f"Invalid slide_index: {slide_index}. Presentation has {len(prs.slides)} slides.",
            }

        # Get slide ID and remove
        slide_id = prs.slides._sldIdLst[slide_index]
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)

        prs.save(file_path)

        logger.info(f"Deleted slide {slide_index} from: {file_path}")

        return {
            "success": True,
            "deleted_index": slide_index,
            "slide_count": len(prs.slides),
            "file_path": file_path,
        }

    except PPTXEditorError as e:
        return {"success": False, "error": str(e)}
    except Exception as e:
        logger.error(f"Failed to delete slide: {e}", exc_info=True)
        return {"success": False, "error": f"Failed to delete slide: {str(e)}"}


@tool_wrapper()
def reorder_slides_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Reorder slides in a presentation.

    Args:
        params: Dictionary containing:
            - file_path (str, required): Path to the PPTX file
            - new_order (list, required): List of slide indices in new order
              e.g., [2, 0, 1] moves slide 2 to first, slide 0 to second, etc.

    Returns:
        Dictionary with:
            - success (bool): Whether operation succeeded
            - new_order (list): The applied order
            - slide_count (int): Total number of slides
            - error (str, optional): Error message if failed
    """
    status.set_callback(params.pop("_status_callback", None))

    file_path = params.get("file_path")
    new_order = params.get("new_order")

    try:
        PPTXEditor._validate_file_path(file_path)

        if not new_order:
            return {"success": False, "error": "new_order parameter is required"}

        if not isinstance(new_order, list):
            return {"success": False, "error": "new_order must be a list of indices"}

        prs = PPTXEditor._get_presentation(file_path)

        if len(new_order) != len(prs.slides):
            return {
                "success": False,
                "error": f"new_order length ({len(new_order)}) must match slide count ({len(prs.slides)})",
            }

        if set(new_order) != set(range(len(prs.slides))):
            return {
                "success": False,
                "error": "new_order must contain each slide index exactly once",
            }

        # Reorder by manipulating the slide ID list
        sldIdLst = prs.slides._sldIdLst
        slide_ids = list(sldIdLst)

        # Clear and re-add in new order
        for slide_id in slide_ids:
            sldIdLst.remove(slide_id)

        for idx in new_order:
            sldIdLst.append(slide_ids[idx])

        prs.save(file_path)

        logger.info(f"Reordered slides in: {file_path}")

        return {
            "success": True,
            "new_order": new_order,
            "slide_count": len(prs.slides),
            "file_path": file_path,
        }

    except PPTXEditorError as e:
        return {"success": False, "error": str(e)}
    except Exception as e:
        logger.error(f"Failed to reorder slides: {e}", exc_info=True)
        return {"success": False, "error": f"Failed to reorder slides: {str(e)}"}


@tool_wrapper()
def add_image_to_slide_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Add an image to a slide.

    Args:
        params: Dictionary containing:
            - file_path (str, required): Path to the PPTX file
            - slide_index (int, required): Index of slide to add image to (0-based)
            - image_path (str, required): Path to the image file
            - position (dict, optional): Position with 'left', 'top', 'width', 'height' in inches
              Default: {'left': 1, 'top': 2, 'width': 6, 'height': 4}

    Returns:
        Dictionary with:
            - success (bool): Whether operation succeeded
            - slide_index (int): Index of slide with image
            - image_path (str): Path to added image
            - error (str, optional): Error message if failed
    """
    status.set_callback(params.pop("_status_callback", None))

    file_path = params.get("file_path")
    slide_index = params.get("slide_index")
    image_path = params.get("image_path")
    position = params.get("position", {})

    try:
        from pptx.util import Inches

        PPTXEditor._validate_file_path(file_path)

        if slide_index is None:
            return {"success": False, "error": "slide_index parameter is required"}

        if not image_path:
            return {"success": False, "error": "image_path parameter is required"}

        if not os.path.exists(image_path):
            return {"success": False, "error": f"Image file not found: {image_path}"}

        prs = PPTXEditor._get_presentation(file_path)

        if slide_index < 0 or slide_index >= len(prs.slides):
            return {
                "success": False,
                "error": f"Invalid slide_index: {slide_index}. Presentation has {len(prs.slides)} slides.",
            }

        slide = prs.slides[slide_index]

        # Position defaults
        left = Inches(position.get("left", 1))
        top = Inches(position.get("top", 2))
        width = Inches(position.get("width", 6)) if "width" in position else None
        height = Inches(position.get("height", 4)) if "height" in position else None

        # Add image
        if width and height:
            slide.shapes.add_picture(image_path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(image_path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(image_path, left, top, height=height)
        else:
            slide.shapes.add_picture(image_path, left, top)

        prs.save(file_path)

        logger.info(f"Added image to slide {slide_index} in: {file_path}")

        return {
            "success": True,
            "slide_index": slide_index,
            "image_path": image_path,
            "file_path": file_path,
        }

    except PPTXEditorError as e:
        return {"success": False, "error": str(e)}
    except Exception as e:
        logger.error(f"Failed to add image: {e}", exc_info=True)
        return {"success": False, "error": f"Failed to add image: {str(e)}"}


@tool_wrapper()
def extract_text_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract all text from a presentation.

    Args:
        params: Dictionary containing:
            - file_path (str, required): Path to the PPTX file

    Returns:
        Dictionary with:
            - success (bool): Whether operation succeeded
            - text (str): All extracted text joined by newlines
            - slides_text (list): List of text per slide
            - total_characters (int): Total character count
            - error (str, optional): Error message if failed
    """
    status.set_callback(params.pop("_status_callback", None))

    file_path = params.get("file_path")

    try:
        PPTXEditor._validate_file_path(file_path)
        prs = PPTXEditor._get_presentation(file_path)

        slides_text = []
        all_text = []

        for idx, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                text = PPTXEditor._extract_shape_text(shape)
                if text:
                    slide_texts.append(text)

            slide_combined = "\n".join(slide_texts)
            slides_text.append({"index": idx, "text": slide_combined})
            all_text.append(slide_combined)

        combined_text = "\n\n---\n\n".join(all_text)

        logger.info(f"Extracted text from {len(prs.slides)} slides in: {file_path}")

        return {
            "success": True,
            "text": combined_text,
            "slides_text": slides_text,
            "slide_count": len(prs.slides),
            "total_characters": len(combined_text),
            "file_path": file_path,
        }

    except PPTXEditorError as e:
        return {"success": False, "error": str(e)}
    except Exception as e:
        logger.error(f"Failed to extract text: {e}", exc_info=True)
        return {"success": False, "error": f"Failed to extract text: {str(e)}"}


# =========================================================================
# PptxGenJS Presentation Engine
# =========================================================================

DEFAULT_THEME = {
    "primary": "1E3A5F",
    "secondary": "0066CC",
    "accent": "059669",
    "info_blue": "E0F2FE",
    "success_green": "D1FAE5",
    "warning_amber": "FEF3C7",
    "danger_red": "FEE2E2",
    "light_gray": "F3F4F6",
    "background": "F8FAFC",
    "surface": "F1F5F9",
    "text_primary": "1F2937",
    "text_secondary": "6B7280",
    "text_muted": "9CA3AF",
    "border": "E5E7EB",
    "divider": "D1D5DB",
    "font_heading": "Calibri",
    "font_body": "Calibri",
    "font_mono": "Consolas",
    "title_size": 40,
    "heading_size": 28,
    "subheading_size": 22,
    "body_size": 16,
    "caption_size": 12,
    "corner_radius": 4,
    "padding": 0.3,
    "margin": 0.5,
}


class PptxGenJSEngine:
    """Engine for PptxGenJS-based presentation creation and PPTX reading."""

    NS = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "dc": "http://purl.org/dc/elements/1.1/",
        "cp": "http://schemas.openxmlformats.org/package/2006/metadata/core-properties",
    }

    @staticmethod
    def _ensure_storage():
        Path(PPTX_STORAGE_PATH).mkdir(parents=True, exist_ok=True)

    @classmethod
    def _unique_path(cls, name, ext=".pptx"):
        cls._ensure_storage()
        safe = "".join(c for c in name if c.isalnum() or c in "._-") or "presentation"
        if not safe.endswith(ext):
            safe += ext
        return os.path.join(PPTX_STORAGE_PATH, f"{uuid.uuid4().hex[:8]}_{safe}")

    @classmethod
    def _script_path(cls, name):
        cls._ensure_storage()
        return os.path.join(PPTX_STORAGE_PATH, f"{uuid.uuid4().hex[:8]}_{name}.mjs")

    @staticmethod
    def _emu_to_inches(emu):
        return int(emu) / 914400.0

    @classmethod
    def _get_text_runs(cls, element):
        """Extract text runs with formatting from XML element."""
        ns = cls.NS
        runs = []
        for r in element.findall(".//a:r", ns):
            text_parts = [t.text for t in r.findall(".//a:t", ns) if t.text]
            text = "".join(text_parts)
            if not text:
                continue
            run = {"text": text}
            rpr = r.find("a:rPr", ns)
            if rpr is not None:
                if rpr.get("b") == "1":
                    run["bold"] = True
                if rpr.get("i") == "1":
                    run["italic"] = True
                srgb = rpr.find(".//a:srgbClr", ns)
                if srgb is not None and srgb.get("val"):
                    run["color"] = srgb.get("val").upper()
                sz = rpr.get("sz")
                if sz:
                    try:
                        run["fontSize"] = int(int(sz) / 100)
                    except (ValueError, TypeError):
                        pass
            runs.append(run)
        return runs

    @classmethod
    def _parse_slide_xml(cls, slide_xml, slide_rels, media_map):
        """Parse a single slide XML into structured data."""
        ns = cls.NS
        data = {"text_elements": [], "images": [], "tables": []}
        try:
            root = ET.fromstring(slide_xml)

            for sp in root.findall(".//p:sp", ns):
                name = ""
                nv = sp.find("p:nvSpPr", ns)
                if nv is not None:
                    cnv = nv.find("p:cNvPr", ns)
                    if cnv is not None:
                        name = cnv.get("name", "")

                x, y, w, h = 0, 0, 0, 0
                xfrm = sp.find(".//a:xfrm", ns)
                if xfrm is not None:
                    off = xfrm.find("a:off", ns)
                    ext = xfrm.find("a:ext", ns)
                    if off is not None:
                        x = cls._emu_to_inches(off.get("x", 0))
                        y = cls._emu_to_inches(off.get("y", 0))
                    if ext is not None:
                        w = cls._emu_to_inches(ext.get("cx", 0))
                        h = cls._emu_to_inches(ext.get("cy", 0))

                tx = sp.find(".//a:txBody", ns)
                if tx is not None:
                    runs = cls._get_text_runs(tx)
                    if runs:
                        full = "".join(r.get("text", "") for r in runs)
                        data["text_elements"].append(
                            {
                                "text": full,
                                "runs": runs,
                                "x": x,
                                "y": y,
                                "w": w,
                                "h": h,
                                "name": name,
                            }
                        )

            for pic in root.findall(".//p:pic", ns):
                blip = pic.find(".//a:blip", ns)
                if blip is not None:
                    eid = blip.get(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                    )
                    if eid and eid in slide_rels and slide_rels[eid] in media_map:
                        x, y, w, h = 0, 0, 0, 0
                        xfrm = pic.find(".//a:xfrm", ns)
                        if xfrm is not None:
                            off = xfrm.find("a:off", ns)
                            ext = xfrm.find("a:ext", ns)
                            if off is not None:
                                x = cls._emu_to_inches(off.get("x", 0))
                                y = cls._emu_to_inches(off.get("y", 0))
                            if ext is not None:
                                w = cls._emu_to_inches(ext.get("cx", 0))
                                h = cls._emu_to_inches(ext.get("cy", 0))
                        data["images"].append(
                            {"path": media_map[slide_rels[eid]], "x": x, "y": y, "w": w, "h": h}
                        )

            for tbl in root.findall(".//a:tbl", ns):
                rows = []
                for tr in tbl.findall(".//a:tr", ns):
                    row = []
                    for tc in tr.findall(".//a:tc", ns):
                        runs = cls._get_text_runs(tc)
                        row.append("".join(r.get("text", "") for r in runs))
                    if row:
                        rows.append(row)
                if rows:
                    data["tables"].append({"rows": rows})

        except Exception as e:
            logger.warning(f"Error parsing slide XML: {e}")
        return data

    @classmethod
    def _slide_to_js(cls, slide_data, idx, theme):
        """Convert parsed slide data to PptxGenJS JavaScript lines."""
        t = theme
        lines = [
            f"// Slide {idx + 1}",
            f"const slide{idx} = pres.addSlide();",
            f"slide{idx}.background = {{ color: BG_COLOR }};",
            "",
        ]

        for i, te in enumerate(slide_data.get("text_elements", [])):
            text = te.get("text", "").strip()
            if not text:
                continue
            x, y, w, h = te.get("x", 0.5), te.get("y", 0.5), te.get("w", 9), te.get("h", 1)
            lines.append(f"slide{idx}.addText({json.dumps(text)}, {{")
            lines.append(f"  x: {x:.2f}, y: {y:.2f}, w: {w:.2f}, h: {h:.2f},")
            lines.append(f"  fontFace: FONT_BODY, fontSize: BODY_SIZE, color: TEXT_PRIMARY")
            lines.append(f"}});")
            lines.append("")

        for i, img in enumerate(slide_data.get("images", [])):
            path = img.get("path", "")
            if not path or not os.path.exists(path):
                continue
            x, y, w, h = img.get("x", 1), img.get("y", 1), img.get("w", 4), img.get("h", 3)
            lines.append(f"slide{idx}.addImage({{")
            lines.append(f"  path: {json.dumps(path)},")
            lines.append(f"  x: {x:.2f}, y: {y:.2f}, w: {w:.2f}, h: {h:.2f}")
            lines.append(f"}});")
            lines.append("")

        for i, tbl in enumerate(slide_data.get("tables", [])):
            rows = tbl.get("rows", [])
            if not rows:
                continue
            table_data = []
            for ri, row in enumerate(rows):
                table_row = []
                for cell in row:
                    opts = {}
                    if ri == 0:
                        opts = {
                            "bold": True,
                            "fill": {"color": t["secondary"]},
                            "color": "FFFFFF",
                            "fontSize": 14,
                            "valign": "middle",
                        }
                    else:
                        opts = {
                            "fill": {
                                "color": (
                                    t["surface"] if ri % 2 == 0 else t.get("light_gray", "F3F4F6")
                                )
                            },
                            "color": t["text_primary"],
                            "fontSize": 13,
                            "valign": "middle",
                        }
                    table_row.append({"text": str(cell), "options": opts})
                table_data.append(table_row)

            lines.append(f"slide{idx}.addTable({json.dumps(table_data)}, {{")
            lines.append(f"  x: 0.50, y: 1.00, w: 9.00,")
            lines.append(f"  border: {{ pt: 0.5, color: '{t.get('border', 'E5E7EB')}' }},")
            lines.append(f"  margin: 0.1, rowH: 0.4")
            lines.append(f"}});")
            lines.append("")

        return lines

    @classmethod
    def _extract_paragraph_md(cls, p_elem):
        """Extract a paragraph element to markdown."""
        ns = cls.NS
        text = ""
        has_bold = False
        for r in p_elem.findall(".//a:r", ns):
            rpr = r.find("a:rPr", ns)
            if rpr is not None and rpr.get("b") == "1":
                has_bold = True
            for t in r.findall("a:t", ns):
                if t.text:
                    text += t.text

        if not text.strip():
            return ""
        text = text.strip()

        ppr = p_elem.find("a:pPr", ns)
        is_bullet = is_numbered = False
        level = 0

        if ppr is not None:
            lvl = ppr.get("lvl")
            if lvl:
                try:
                    level = int(lvl)
                except ValueError:
                    pass
            if ppr.find("a:buAutoNum", ns) is not None:
                is_numbered = True
            elif ppr.find("a:buNone", ns) is None and (
                ppr.find("a:buChar", ns) is not None or ppr.find("a:buFont", ns) is not None
            ):
                is_bullet = True

        if not is_bullet and not is_numbered and text[0] in "•–▪►‣◦○":
            is_bullet = True
            text = text.lstrip("•–▪►‣◦○ ")

        indent = "  " * level
        if is_numbered:
            return f"{indent}1. {text}"
        elif is_bullet:
            return f"{indent}- {text}"
        elif has_bold and len(text) < 80:
            return f"{indent}**{text}**"
        return f"{indent}{text}"

    @classmethod
    def write_script(cls, script_content, filename="custom_pptx.mjs"):
        """Save an LLM-generated PptxGenJS Node.js script to disk."""
        try:
            path = cls._script_path(filename.replace(".mjs", ""))
            with open(path, "w") as f:
                f.write(script_content)
            return tool_response(
                script_path=path,
                script_size_bytes=len(script_content.encode("utf-8")),
                next_step=f"Call execute_pptx_script_tool with script_path='{path}'",
            )
        except Exception as e:
            return tool_error(str(e))

    @classmethod
    def execute_script(cls, script_path, timeout=60):
        """Execute a PptxGenJS script via Node.js."""
        if not os.path.exists(script_path):
            return tool_error(f"Script not found: {script_path}")

        script_dir = os.path.dirname(script_path)
        script_name = os.path.basename(script_path)

        try:
            subprocess.run(
                ["npm", "install", "pptxgenjs"],
                cwd=script_dir,
                capture_output=True,
                text=True,
                timeout=30,
            )

            result = subprocess.run(
                ["node", script_name],
                cwd=script_dir,
                capture_output=True,
                text=True,
                timeout=timeout,
            )

            output = result.stdout
            pptx_path = None
            for line in output.split("\n"):
                if "saved to:" in line.lower():
                    parts = line.split("saved to:")
                    if len(parts) > 1:
                        pptx_path = parts[1].strip().strip("'\"")
                        break

            if result.returncode == 0:
                return tool_response(output=output, return_code=0, pptx_path=pptx_path)
            else:
                return tool_error(
                    f"Script failed (code {result.returncode}): {result.stderr or output}"
                )

        except subprocess.TimeoutExpired:
            return tool_error(f"Script timed out after {timeout}s")
        except FileNotFoundError:
            return tool_error("Node.js not found. Ensure Node.js is installed.")
        except Exception as e:
            return tool_error(str(e))

    @classmethod
    def read_to_markdown(cls, pptx_path, include_notes=True, include_metadata=True):
        """Read PPTX file and convert to Markdown."""
        if not os.path.exists(pptx_path):
            return tool_error(f"File not found: {pptx_path}")
        if not pptx_path.lower().endswith(".pptx"):
            return tool_error(f"Not a PPTX file: {pptx_path}")

        ns = cls.NS
        try:
            md_parts, metadata = [], {}
            slide_count = total_tables = total_images = 0

            with zipfile.ZipFile(pptx_path, "r") as z:
                all_files = z.namelist()

                if include_metadata and "docProps/core.xml" in all_files:
                    try:
                        root = ET.fromstring(
                            z.read("docProps/core.xml").decode("utf-8", errors="replace")
                        )
                        for tag, key in [
                            (".//dc:title", "title"),
                            (".//dc:creator", "author"),
                            (".//dc:subject", "subject"),
                        ]:
                            el = root.find(tag, ns)
                            if el is not None and el.text:
                                metadata[key] = el.text
                    except Exception:
                        pass

                md_parts.append(f"# {metadata.get('title', Path(pptx_path).stem)}\n")
                if metadata.get("author"):
                    md_parts.append(f"*Author: {metadata['author']}*\n")
                md_parts.append("---\n")

                slide_files = sorted(
                    [
                        f
                        for f in all_files
                        if f.startswith("ppt/slides/slide")
                        and f.endswith(".xml")
                        and ".rels" not in f
                    ],
                    key=lambda x: int("".join(filter(str.isdigit, x.split("/")[-1])) or "0"),
                )

                for slide_file in slide_files:
                    slide_count += 1
                    root = ET.fromstring(z.read(slide_file).decode("utf-8", errors="replace"))

                    title_text = subtitle_text = ""
                    body_parts, slide_images, slide_tables = [], [], []

                    # Read slide relationships for images
                    slide_rels = {}
                    rels_path = slide_file.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
                    if rels_path in all_files:
                        try:
                            rr = ET.fromstring(z.read(rels_path).decode("utf-8", errors="replace"))
                            rns = {
                                "r": "http://schemas.openxmlformats.org/package/2006/relationships"
                            }
                            for rel in rr.findall(".//r:Relationship", rns):
                                rid = rel.get("Id")
                                target = rel.get("Target", "")
                                if rid and ("media" in target or "image" in target.lower()):
                                    slide_rels[rid] = os.path.basename(target)
                        except Exception:
                            pass

                    # Extract text from shapes
                    for sp in root.findall(".//p:sp", ns):
                        ph_type = ""
                        nvPr = sp.find(".//p:nvPr", ns)
                        if nvPr is not None:
                            ph = nvPr.find(".//p:ph", ns)
                            if ph is not None:
                                ph_type = ph.get("type", "")

                        is_title = ph_type in ("title", "ctrTitle")

                        tx = sp.find(".//a:txBody", ns)
                        if tx is not None:
                            lines = [cls._extract_paragraph_md(p) for p in tx.findall("a:p", ns)]
                            full = "\n".join(l for l in lines if l)
                            if full.strip():
                                if is_title:
                                    title_text = full.strip()
                                else:
                                    body_parts.append(full)

                    # Tables
                    for tbl in root.findall(".//a:tbl", ns):
                        total_tables += 1
                        rows = []
                        for tr in tbl.findall(".//a:tr", ns):
                            row = []
                            for tc in tr.findall(".//a:tc", ns):
                                cell = "".join(
                                    t.text for t in tc.findall(".//a:t", ns) if t.text
                                ).strip()
                                row.append(cell)
                            if row:
                                rows.append(row)
                        if rows:
                            header = rows[0]
                            tmd = [
                                "| " + " | ".join(header) + " |",
                                "|" + "|".join(["---"] * len(header)) + "|",
                            ]
                            for dr in rows[1:]:
                                padded = dr + [""] * (len(header) - len(dr))
                                tmd.append("| " + " | ".join(padded[: len(header)]) + " |")
                            slide_tables.append("\n".join(tmd))

                    # Images
                    for pic in root.findall(".//p:pic", ns):
                        blip = pic.find(".//a:blip", ns)
                        if blip is not None:
                            eid = blip.get(
                                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                            )
                            if eid and eid in slide_rels:
                                total_images += 1
                                slide_images.append(slide_rels[eid])

                    # Build slide markdown
                    header = f"## Slide {slide_count}"
                    if title_text:
                        header += f": {title_text}"
                    md_parts.append(f"\n{header}\n")
                    for part in body_parts:
                        md_parts.append(part)
                    for tmd in slide_tables:
                        md_parts.append(f"\n{tmd}\n")
                    if slide_images:
                        md_parts.append(f"\n*Images: {', '.join(f'`{i}`' for i in slide_images)}*")

                    # Notes
                    if include_notes:
                        snum = slide_file.split("/")[-1].replace("slide", "").replace(".xml", "")
                        nf = f"ppt/notesSlides/notesSlide{snum}.xml"
                        if nf in all_files:
                            try:
                                nr = ET.fromstring(z.read(nf).decode("utf-8", errors="replace"))
                                notes = " ".join(
                                    t.text.strip()
                                    for t in nr.findall(".//a:t", ns)
                                    if t.text and t.text.strip()
                                )
                                if notes and notes.lower() not in (
                                    "click to add notes",
                                    "notes",
                                    "click to add text",
                                ):
                                    md_parts.append(f"\n> **Speaker Notes:** {notes}")
                            except Exception:
                                pass
                    md_parts.append("")

            return tool_response(
                markdown="\n".join(md_parts),
                slides_count=slide_count,
                tables_count=total_tables,
                images_count=total_images,
                metadata=metadata,
            )

        except zipfile.BadZipFile:
            return tool_error(f"Invalid PPTX file: {pptx_path}")
        except Exception as e:
            return tool_error(str(e))

    @classmethod
    def read_to_js(cls, pptx_path, output_path=None, preserve_images=True):
        """Read PPTX and convert to PptxGenJS JavaScript code."""
        if not os.path.exists(pptx_path):
            return tool_error(f"File not found: {pptx_path}")
        if not pptx_path.lower().endswith(".pptx"):
            return tool_error(f"Not a PPTX file: {pptx_path}")

        warnings = []
        slides_data = []
        metadata = {}
        media_map = {}

        try:
            media_dir = os.path.join(PPTX_STORAGE_PATH, f"pptx_media_{uuid.uuid4().hex[:8]}")
            cls._ensure_storage()
            os.makedirs(media_dir, exist_ok=True)

            with zipfile.ZipFile(pptx_path, "r") as z:
                # Get slide list from presentation.xml
                slide_ids = []
                slide_id_to_file = {}
                try:
                    pres_xml = z.read("ppt/presentation.xml")
                    pres_root = ET.fromstring(pres_xml)
                    ns_p = {"p": cls.NS["p"]}
                    for sid in pres_root.findall(".//p:sldId", ns_p):
                        rid = sid.get(
                            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
                        )
                        if rid:
                            slide_ids.append(rid)

                    pres_rels = ET.fromstring(z.read("ppt/_rels/presentation.xml.rels"))
                    rns = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}
                    for rel in pres_rels.findall(".//r:Relationship", rns):
                        rid = rel.get("Id")
                        if rid in slide_ids:
                            slide_id_to_file[rid] = rel.get("Target")
                except Exception as e:
                    warnings.append(str(e))

                # Metadata
                try:
                    core = ET.fromstring(z.read("docProps/core.xml"))
                    for tag, key in [(".//dc:title", "title"), (".//dc:creator", "author")]:
                        el = core.find(tag, cls.NS)
                        if el is not None and el.text:
                            metadata[key] = el.text
                except Exception:
                    pass

                # Theme colors
                theme_colors = {}
                try:
                    tfs = [
                        f
                        for f in z.namelist()
                        if f.startswith("ppt/theme/theme") and f.endswith(".xml")
                    ]
                    if tfs:
                        troot = ET.fromstring(z.read(tfs[0]))
                        ns_a = {"a": cls.NS["a"]}
                        cs = troot.find(".//a:clrScheme", ns_a)
                        if cs is not None:
                            for child in cs:
                                tag = child.tag.split("}")[-1]
                                srgb = child.find(".//a:srgbClr", ns_a)
                                if srgb is not None and srgb.get("val"):
                                    theme_colors[tag] = srgb.get("val").upper()
                except Exception:
                    pass

                # Extract media
                if preserve_images:
                    for mf in z.namelist():
                        if mf.startswith("ppt/media/") and os.path.splitext(mf)[1].lower() in [
                            ".png",
                            ".jpg",
                            ".jpeg",
                            ".gif",
                            ".svg",
                        ]:
                            ep = os.path.join(media_dir, os.path.basename(mf))
                            with open(ep, "wb") as f:
                                f.write(z.read(mf))
                            media_map[mf] = ep

                # Parse slides
                for si, sid in enumerate(slide_ids):
                    sf = slide_id_to_file.get(sid)
                    if not sf:
                        continue
                    if not sf.startswith("ppt/"):
                        sf = "ppt/" + sf.lstrip("/")
                    try:
                        slide_xml = z.read(sf)
                        slide_rels = {}
                        rf = sf.replace("ppt/slides/", "ppt/slides/_rels/") + ".rels"
                        if rf in z.namelist():
                            rr = ET.fromstring(z.read(rf))
                            rns2 = {
                                "r": "http://schemas.openxmlformats.org/package/2006/relationships"
                            }
                            for rel in rr.findall(".//r:Relationship", rns2):
                                rid = rel.get("Id")
                                tgt = rel.get("Target")
                                if rid and tgt:
                                    if tgt.startswith("../"):
                                        tgt = "ppt/" + tgt[3:]
                                    elif not tgt.startswith("ppt/"):
                                        tgt = os.path.dirname(sf) + "/" + tgt
                                    slide_rels[rid] = tgt

                        slides_data.append(cls._parse_slide_xml(slide_xml, slide_rels, media_map))
                    except Exception as e:
                        warnings.append(f"Slide {si + 1}: {e}")
                        slides_data.append({"text_elements": [], "images": [], "tables": []})

            # Build effective theme
            et = dict(DEFAULT_THEME)
            mapping = {
                "dk1": "text_primary",
                "dk2": "primary",
                "lt1": "background",
                "lt2": "surface",
                "accent1": "secondary",
                "accent2": "accent",
            }
            for tk, ok in mapping.items():
                if tk in theme_colors:
                    et[ok] = theme_colors[tk]

            # Generate JS
            script_lines = [
                f"import pptxgen from 'pptxgenjs';",
                f"const pres = new pptxgen();",
                f"pres.title = {json.dumps(metadata.get('title', 'Presentation'))};",
                f"pres.author = {json.dumps(metadata.get('author', 'Jotty AI'))};",
                "",
                f"const PRIMARY = '{et['primary']}';",
                f"const SECONDARY = '{et['secondary']}';",
                f"const BG_COLOR = '{et['background']}';",
                f"const TEXT_PRIMARY = '{et['text_primary']}';",
                f"const TEXT_SECONDARY = '{et['text_secondary']}';",
                f"const FONT_HEADING = '{et['font_heading']}';",
                f"const FONT_BODY = '{et['font_body']}';",
                f"const TITLE_SIZE = {et['title_size']};",
                f"const HEADING_SIZE = {et['heading_size']};",
                f"const BODY_SIZE = {et['body_size']};",
                "",
            ]

            for si, sd in enumerate(slides_data):
                script_lines.extend(cls._slide_to_js(sd, si, et))
                script_lines.append("")

            script_lines.append("// await pres.writeFile({ fileName: 'output.pptx' });")

            script = "\n".join(script_lines)
            if output_path is None:
                output_path = cls._script_path("read_pptx")
            with open(output_path, "w") as f:
                f.write(script)

            return tool_response(
                script=script,
                script_path=output_path,
                slides_extracted=len(slides_data),
                images_extracted=len(media_map),
                metadata=metadata,
                warnings=warnings,
            )

        except zipfile.BadZipFile:
            return tool_error(f"Invalid PPTX file: {pptx_path}")
        except Exception as e:
            return tool_error(str(e))

    @staticmethod
    def get_api_reference():
        """Get PptxGenJS API reference for agents."""
        return tool_response(
            version="4.0.1",
            docs="https://gitbrent.github.io/PptxGenJS/",
            charts={
                "types": [
                    "area",
                    "bar",
                    "bar3D",
                    "bubble",
                    "doughnut",
                    "line",
                    "pie",
                    "radar",
                    "scatter",
                ],
                "data_format": [{"name": "Series", "labels": ["A", "B"], "values": [100, 200]}],
                "options": [
                    "x",
                    "y",
                    "w",
                    "h",
                    "chartColors",
                    "showLegend",
                    "legendPos",
                    "showLabel",
                    "showValue",
                    "showPercent",
                    "lineSmooth",
                    "holeSize",
                ],
            },
            text={
                "methods": ["addText(text, opts)", "addText([{text, options}], opts)"],
                "options": [
                    "fontFace",
                    "fontSize",
                    "color",
                    "bold",
                    "italic",
                    "align",
                    "valign",
                    "bullet",
                    "paraSpaceAfter",
                    "lineSpacing",
                    "shadow",
                    "rotate",
                ],
            },
            tables={
                "methods": ["addTable(rows, opts)"],
                "options": ["border", "fill", "colW", "rowH", "margin", "colspan", "rowspan"],
            },
            shapes=[
                "rect",
                "roundRect",
                "ellipse",
                "triangle",
                "line",
                "arrow",
                "chevron",
                "pentagon",
                "hexagon",
                "star5",
                "cloud",
                "diamond",
            ],
            images={"methods": ["addImage(opts)"], "sources": ["path", "data (base64)", "url"]},
            default_theme=DEFAULT_THEME,
        )

    @staticmethod
    def get_chart_template(chart_type):
        """Get chart data template for a given chart type."""
        templates = {
            "bar": {
                "data": [{"name": "Q1", "labels": ["A", "B", "C"], "values": [4500, 2500, 1800]}],
                "options": {"barDir": "bar", "showValue": True},
            },
            "line": {
                "data": [
                    {
                        "name": "Revenue",
                        "labels": ["Jan", "Feb", "Mar", "Apr"],
                        "values": [1000, 1200, 1100, 1400],
                    }
                ],
                "options": {"lineSmooth": True, "lineDataSymbol": "circle"},
            },
            "pie": {
                "data": [
                    {
                        "name": "Share",
                        "labels": ["A", "B", "C", "Other"],
                        "values": [35, 25, 20, 20],
                    }
                ],
                "options": {"showPercent": True},
            },
            "doughnut": {
                "data": [
                    {
                        "name": "Budget",
                        "labels": ["R&D", "Marketing", "Ops"],
                        "values": [40, 35, 25],
                    }
                ],
                "options": {"holeSize": 50, "showPercent": True},
            },
            "area": {
                "data": [
                    {
                        "name": "Traffic",
                        "labels": ["W1", "W2", "W3", "W4"],
                        "values": [5000, 7500, 6000, 8500],
                    }
                ],
                "options": {"chartColorsOpacity": 50},
            },
            "radar": {
                "data": [
                    {
                        "name": "Team A",
                        "labels": ["Speed", "Power", "Skill", "Stamina"],
                        "values": [80, 70, 90, 75],
                    },
                    {
                        "name": "Team B",
                        "labels": ["Speed", "Power", "Skill", "Stamina"],
                        "values": [75, 85, 70, 80],
                    },
                ],
                "options": {"radarStyle": "standard"},
            },
        }
        return tool_response(
            template=templates.get(chart_type, templates["bar"]), chart_type=chart_type
        )

    @staticmethod
    def search_docs(query):
        """Search PptxGenJS docs for a feature."""
        doc_map = {
            "chart": {
                "url": "https://gitbrent.github.io/PptxGenJS/docs/api-charts/",
                "topics": ["bar", "line", "pie", "doughnut", "area", "radar", "scatter"],
            },
            "table": {
                "url": "https://gitbrent.github.io/PptxGenJS/docs/api-tables/",
                "topics": ["border", "cell", "rowspan", "colspan", "header", "fill"],
            },
            "text": {
                "url": "https://gitbrent.github.io/PptxGenJS/docs/api-text/",
                "topics": ["font", "bullet", "align", "color", "bold"],
            },
            "image": {
                "url": "https://gitbrent.github.io/PptxGenJS/docs/api-images/",
                "topics": ["path", "data", "base64", "sizing"],
            },
            "shape": {
                "url": "https://gitbrent.github.io/PptxGenJS/docs/api-shapes/",
                "topics": ["rect", "ellipse", "line", "arrow", "fill", "shadow"],
            },
        }
        q = query.lower()
        results = []
        for cat, info in doc_map.items():
            if cat in q or any(t in q for t in info["topics"]):
                results.append({"category": cat, "url": info["url"]})
        if not results:
            results = [{"category": "general", "url": "https://gitbrent.github.io/PptxGenJS/"}]
        return tool_response(results=results, count=len(results))

    @staticmethod
    def save_file(source_path, destination_path):
        """Copy a generated file to a user-specified location."""
        if not os.path.exists(source_path):
            return tool_error(f"Source not found: {source_path}")
        try:
            dest = os.path.expanduser(os.path.expandvars(destination_path))
            dest_dir = os.path.dirname(dest)
            if dest_dir:
                os.makedirs(dest_dir, exist_ok=True)
            shutil.copy2(source_path, dest)
            return tool_response(
                saved_to=os.path.abspath(dest), file_size_bytes=os.path.getsize(dest)
            )
        except PermissionError:
            return tool_error(f"Permission denied: {destination_path}")
        except Exception as e:
            return tool_error(str(e))


# =========================================================================
# PptxGenJS Tool Functions
# =========================================================================


@tool_wrapper(required_params=["script_content"])
def write_pptx_script_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Save an LLM-generated PptxGenJS Node.js script to disk.

    Args:
        params: Dictionary containing:
            - script_content (str, required): Full PptxGenJS Node.js script
            - filename (str, optional): Script filename (default: custom_pptx.mjs)

    Returns:
        Dictionary with success, script_path, script_size_bytes, next_step
    """
    status.set_callback(params.pop("_status_callback", None))
    status.emit("Writing", "Saving PptxGenJS script")
    return PptxGenJSEngine.write_script(
        params["script_content"], params.get("filename", "custom_pptx.mjs")
    )


@tool_wrapper(required_params=["script_path"])
def execute_pptx_script_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Execute a PptxGenJS script via Node.js.

    Args:
        params: Dictionary containing:
            - script_path (str, required): Path to .mjs script
            - timeout (int, optional): Max seconds (default: 60)

    Returns:
        Dictionary with success, output, pptx_path
    """
    status.set_callback(params.pop("_status_callback", None))
    status.emit("Executing", "Running PptxGenJS script")
    return PptxGenJSEngine.execute_script(params["script_path"], params.get("timeout", 60))


@tool_wrapper(required_params=["pptx_path"])
def read_pptx_to_markdown_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Read a PPTX file and convert to Markdown.

    Args:
        params: Dictionary containing:
            - pptx_path (str, required): Path to .pptx file
            - include_notes (bool, optional): Include speaker notes (default: True)
            - include_metadata (bool, optional): Include file metadata (default: True)

    Returns:
        Dictionary with success, markdown, slides_count, tables_count, images_count
    """
    status.set_callback(params.pop("_status_callback", None))
    status.emit("Reading", "Converting PPTX to Markdown")
    return PptxGenJSEngine.read_to_markdown(
        params["pptx_path"], params.get("include_notes", True), params.get("include_metadata", True)
    )


@tool_wrapper(required_params=["pptx_path"])
def read_pptx_to_js_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Read a PPTX file and convert to PptxGenJS JavaScript code.

    Args:
        params: Dictionary containing:
            - pptx_path (str, required): Path to .pptx file
            - output_path (str, optional): Path for generated script
            - preserve_images (bool, optional): Extract images (default: True)

    Returns:
        Dictionary with success, script, script_path, slides_extracted
    """
    status.set_callback(params.pop("_status_callback", None))
    status.emit("Converting", "Converting PPTX to PptxGenJS")
    return PptxGenJSEngine.read_to_js(
        params["pptx_path"], params.get("output_path"), params.get("preserve_images", True)
    )


@tool_wrapper()
def get_pptxgenjs_api_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get PptxGenJS API reference for agent exploration.

    Returns:
        Dictionary with charts, text, tables, shapes, images API reference
    """
    status.set_callback(params.pop("_status_callback", None))
    return PptxGenJSEngine.get_api_reference()


@tool_wrapper(required_params=["chart_type"])
def get_chart_data_template_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Get chart data template for a given chart type.

    Args:
        params: Dictionary containing:
            - chart_type (str, required): bar, line, pie, doughnut, area, radar

    Returns:
        Dictionary with template data and options
    """
    status.set_callback(params.pop("_status_callback", None))
    return PptxGenJSEngine.get_chart_template(params["chart_type"])


@tool_wrapper(required_params=["query"])
def search_pptxgenjs_docs_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Search PptxGenJS documentation for a feature.

    Args:
        params: Dictionary containing:
            - query (str, required): Search query (e.g., 'radar chart', 'table border')

    Returns:
        Dictionary with matching doc sections and URLs
    """
    status.set_callback(params.pop("_status_callback", None))
    return PptxGenJSEngine.search_docs(params["query"])


@tool_wrapper(required_params=["source_path", "destination_path"])
def save_file_to_path_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Copy a generated file to a user-specified location.

    Args:
        params: Dictionary containing:
            - source_path (str, required): Path to source file
            - destination_path (str, required): Where to save

    Returns:
        Dictionary with success, saved_to, file_size_bytes
    """
    status.set_callback(params.pop("_status_callback", None))
    return PptxGenJSEngine.save_file(params["source_path"], params["destination_path"])


__all__ = [
    # Python-pptx editor tools
    "read_pptx_tool",
    "get_slide_layouts_tool",
    "add_slide_tool",
    "update_slide_tool",
    "delete_slide_tool",
    "reorder_slides_tool",
    "add_image_to_slide_tool",
    "extract_text_tool",
    # PptxGenJS tools
    "write_pptx_script_tool",
    "execute_pptx_script_tool",
    "read_pptx_to_markdown_tool",
    "read_pptx_to_js_tool",
    "get_pptxgenjs_api_tool",
    "get_chart_data_template_tool",
    "search_pptxgenjs_docs_tool",
    "save_file_to_path_tool",
    # Constants
    "DEFAULT_THEME",
]
