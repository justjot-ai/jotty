"""
Slide Generator Skill

Generates actual PowerPoint (.pptx) slides using python-pptx.
"""
import os
import logging
from typing import Dict, Any, List
from datetime import datetime

logger = logging.getLogger(__name__)


def generate_slides_pdf_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Generate PDF slides directly using reportlab (slide-style pages).

    Args:
        params: Dictionary containing:
            - title (str, required): Presentation title
            - slides (list, required): List of slide dicts with 'title' and 'bullets'
            - subtitle (str, optional): Subtitle for title slide
            - output_path (str, optional): Output directory
            - template (str, optional): Color theme - 'dark', 'light', 'blue'

    Returns:
        Dictionary with success, file_path, slide_count
    """
    try:
        from reportlab.lib.pagesizes import landscape, A4
        from reportlab.lib.colors import HexColor
        from reportlab.pdfgen import canvas
        from reportlab.lib.units import inch
    except ImportError:
        return {'success': False, 'error': 'reportlab not installed. Install with: pip install reportlab'}

    title = params.get('title')
    slides_data = params.get('slides', [])

    if not title or not slides_data:
        return {'success': False, 'error': 'title and slides parameters are required'}

    subtitle = params.get('subtitle', '')
    template = params.get('template', 'dark')
    output_dir = params.get('output_path', os.path.expanduser('~/jotty/presentations'))

    os.makedirs(output_dir, exist_ok=True)

    # Color themes
    themes = {
        'dark': {'bg': '#1e1e1e', 'title': '#ffffff', 'text': '#dcdcdc', 'accent': '#0096ff'},
        'light': {'bg': '#ffffff', 'title': '#1e1e1e', 'text': '#3c3c3c', 'accent': '#0078c8'},
        'blue': {'bg': '#003366', 'title': '#ffffff', 'text': '#c8dcff', 'accent': '#64c8ff'}
    }
    theme = themes.get(template, themes['dark'])

    import textwrap

    def wrap_pdf_text(text, max_width_chars):
        """Wrap text for PDF slides."""
        return textwrap.wrap(text, width=max_width_chars)

    try:
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        safe_title = "".join(c if c.isalnum() or c in ' -_' else '' for c in title)[:30]
        filename = f"{safe_title.replace(' ', '_')}_{timestamp}.pdf"
        file_path = os.path.join(output_dir, filename)

        page_width, page_height = landscape(A4)
        c = canvas.Canvas(file_path, pagesize=landscape(A4))

        def draw_background():
            c.setFillColor(HexColor(theme['bg']))
            c.rect(0, 0, page_width, page_height, fill=True, stroke=False)

        # Title slide
        draw_background()
        c.setFillColor(HexColor(theme['title']))

        # Wrap title if too long
        title_lines = wrap_pdf_text(title, 40)
        title_font_size = 42 if len(title_lines) > 1 else 48
        c.setFont("Helvetica-Bold", title_font_size)

        y_start = page_height / 2 + (len(title_lines) - 1) * 30
        for i, line in enumerate(title_lines[:3]):
            c.drawCentredString(page_width / 2, y_start - i * 55, line)

        if subtitle:
            c.setFillColor(HexColor(theme['text']))
            c.setFont("Helvetica", 22)
            subtitle_lines = wrap_pdf_text(subtitle, 60)
            sub_y = page_height / 2 - 80
            for i, line in enumerate(subtitle_lines[:2]):
                c.drawCentredString(page_width / 2, sub_y - i * 30, line)

        c.showPage()

        # Content slides
        for slide_data in slides_data:
            draw_background()

            slide_title = slide_data.get('title', 'Untitled')
            bullets = slide_data.get('bullets', [])

            # Slide title - with wrapping
            c.setFillColor(HexColor(theme['accent']))
            title_lines = wrap_pdf_text(slide_title, 50)
            title_font = 32 if len(title_lines) > 1 else 36
            c.setFont("Helvetica-Bold", title_font)

            title_y = page_height - 0.7 * inch
            for i, line in enumerate(title_lines[:2]):
                c.drawString(0.5 * inch, title_y - i * 40, line)

            # Accent line
            line_y = title_y - len(title_lines[:2]) * 40 - 10
            c.setStrokeColor(HexColor(theme['accent']))
            c.setLineWidth(3)
            c.line(0.5 * inch, line_y, 3.5 * inch, line_y)

            # Bullets - with dynamic sizing and wrapping
            num_bullets = len(bullets)
            if num_bullets <= 4:
                font_size = 22
                line_spacing = 0.55 * inch
                max_chars = 70
            elif num_bullets <= 6:
                font_size = 20
                line_spacing = 0.48 * inch
                max_chars = 75
            else:
                font_size = 18
                line_spacing = 0.42 * inch
                max_chars = 80

            c.setFillColor(HexColor(theme['text']))
            c.setFont("Helvetica", font_size)
            y_pos = line_y - 0.5 * inch

            for bullet in bullets[:8]:  # Max 8 bullets
                # Wrap long bullets
                bullet_text = bullet if len(bullet) <= max_chars else bullet[:max_chars-3] + "..."
                bullet_lines = wrap_pdf_text(bullet_text, max_chars)

                for j, bline in enumerate(bullet_lines[:2]):  # Max 2 lines per bullet
                    prefix = "•  " if j == 0 else "   "
                    c.drawString(0.7 * inch, y_pos, f"{prefix}{bline}")
                    y_pos -= font_size + 4

                y_pos -= 8  # Extra space between bullets

                if y_pos < 0.8 * inch:  # Stop if running out of space
                    break

            c.showPage()

        c.save()
        logger.info(f"PDF slides generated: {file_path}")

        return {
            'success': True,
            'file_path': file_path,
            'slide_count': len(slides_data) + 1
        }

    except Exception as e:
        logger.error(f"PDF slide generation failed: {e}", exc_info=True)
        return {'success': False, 'error': f'PDF generation failed: {str(e)}'}


def generate_slides_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Generate PowerPoint slides from structured content.

    Args:
        params: Dictionary containing:
            - title (str, required): Presentation title
            - slides (list, required): List of slide dicts with 'title' and 'bullets'
            - subtitle (str, optional): Subtitle for title slide
            - author (str, optional): Author name
            - output_path (str, optional): Output file path
            - template (str, optional): Color theme - 'dark', 'light', 'blue' (default: 'dark')

    Returns:
        Dictionary with:
            - success (bool): Whether generation succeeded
            - file_path (str): Path to generated PPTX file
            - slide_count (int): Number of slides created
            - error (str, optional): Error message if failed
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        return {
            'success': False,
            'error': 'python-pptx not installed. Install with: pip install python-pptx'
        }

    title = params.get('title')
    slides_data = params.get('slides', [])

    if not title:
        return {'success': False, 'error': 'title parameter is required'}

    if not slides_data:
        return {'success': False, 'error': 'slides parameter is required (list of slide dicts)'}

    subtitle = params.get('subtitle', '')
    author = params.get('author', 'Jotty Slide Generator')
    template = params.get('template', 'dark')
    output_dir = params.get('output_path', os.path.expanduser('~/jotty/presentations'))

    os.makedirs(output_dir, exist_ok=True)

    # Color themes
    themes = {
        'dark': {
            'bg': RGBColor(30, 30, 30),
            'title_color': RGBColor(255, 255, 255),
            'text_color': RGBColor(220, 220, 220),
            'accent': RGBColor(0, 150, 255)
        },
        'light': {
            'bg': RGBColor(255, 255, 255),
            'title_color': RGBColor(30, 30, 30),
            'text_color': RGBColor(60, 60, 60),
            'accent': RGBColor(0, 120, 200)
        },
        'blue': {
            'bg': RGBColor(0, 51, 102),
            'title_color': RGBColor(255, 255, 255),
            'text_color': RGBColor(200, 220, 255),
            'accent': RGBColor(100, 200, 255)
        }
    }

    theme = themes.get(template, themes['dark'])

    def wrap_text(text, max_chars=80):
        """Wrap long text to multiple lines."""
        if len(text) <= max_chars:
            return text
        words = text.split()
        lines = []
        current_line = []
        current_len = 0
        for word in words:
            if current_len + len(word) + 1 <= max_chars:
                current_line.append(word)
                current_len += len(word) + 1
            else:
                if current_line:
                    lines.append(' '.join(current_line))
                current_line = [word]
                current_len = len(word)
        if current_line:
            lines.append(' '.join(current_line))
        return '\n'.join(lines)

    try:
        from pptx.enum.text import MSO_AUTO_SIZE

        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        prs.slide_height = Inches(7.5)

        # Title slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Set background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = theme['bg']
        background.line.fill.background()

        # Title text - with word wrap
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = wrap_text(title, 50)
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = theme['title_color']
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
            sub_frame = sub_box.text_frame
            sub_frame.word_wrap = True
            sub_para = sub_frame.paragraphs[0]
            sub_para.text = wrap_text(subtitle, 70)
            sub_para.font.size = Pt(24)
            sub_para.font.color.rgb = theme['text_color']
            sub_para.alignment = PP_ALIGN.CENTER

        # Content slides
        for slide_data in slides_data:
            slide_title = slide_data.get('title', 'Untitled')
            bullets = slide_data.get('bullets', [])

            slide = prs.slides.add_slide(slide_layout)

            # Background
            background = slide.shapes.add_shape(
                1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
            )
            background.fill.solid()
            background.fill.fore_color.rgb = theme['bg']
            background.line.fill.background()

            # Slide title - with word wrap
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1.2))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_para.text = wrap_text(slide_title, 60)
            title_para.font.size = Pt(32)
            title_para.font.bold = True
            title_para.font.color.rgb = theme['accent']

            # Accent line under title
            line = slide.shapes.add_shape(
                1, Inches(0.5), Inches(1.4), Inches(3), Inches(0.04)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = theme['accent']
            line.line.fill.background()

            # Bullets - with proper sizing based on count
            if bullets:
                # Adjust font size based on number of bullets and content length
                max_bullet_len = max(len(b) for b in bullets) if bullets else 0
                num_bullets = len(bullets)

                # Dynamic font sizing
                if num_bullets <= 3 and max_bullet_len < 80:
                    font_size = 24
                    spacing = 20
                elif num_bullets <= 5 and max_bullet_len < 100:
                    font_size = 20
                    spacing = 16
                else:
                    font_size = 18
                    spacing = 12

                content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12), Inches(5.5))
                text_frame = content_box.text_frame
                text_frame.word_wrap = True

                for i, bullet in enumerate(bullets[:8]):  # Max 8 bullets per slide
                    if i == 0:
                        para = text_frame.paragraphs[0]
                    else:
                        para = text_frame.add_paragraph()

                    # Truncate very long bullets
                    bullet_text = bullet if len(bullet) <= 150 else bullet[:147] + "..."
                    para.text = f"•  {bullet_text}"
                    para.font.size = Pt(font_size)
                    para.font.color.rgb = theme['text_color']
                    para.space_after = Pt(spacing)
                    para.line_spacing = 1.2

        # Save
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        safe_title = "".join(c if c.isalnum() or c in ' -_' else '' for c in title)[:30]
        filename = f"{safe_title.replace(' ', '_')}_{timestamp}.pptx"
        file_path = os.path.join(output_dir, filename)

        prs.save(file_path)

        logger.info(f"Slides generated: {file_path}")

        return {
            'success': True,
            'file_path': file_path,
            'slide_count': len(slides_data) + 1,  # +1 for title slide
            'title': title
        }

    except Exception as e:
        logger.error(f"Slide generation failed: {e}", exc_info=True)
        return {
            'success': False,
            'error': f'Slide generation failed: {str(e)}'
        }


async def generate_slides_from_topic_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Generate PowerPoint slides from a topic using AI.

    Args:
        params: Dictionary containing:
            - topic (str, required): Topic to create slides about
            - n_slides (int, optional): Number of content slides (default: 10)
            - template (str, optional): Color theme - 'dark', 'light', 'blue'
            - output_path (str, optional): Output directory
            - export_as (str, optional): 'pptx', 'pdf', or 'both' (default: 'pptx')
            - send_telegram (bool, optional): Send to Telegram (default: True)

    Returns:
        Dictionary with:
            - success (bool): Whether generation succeeded
            - file_path (str): Path to generated PPTX file
            - pdf_path (str): Path to PDF version (if export_as includes pdf)
            - slide_count (int): Number of slides
            - telegram_sent (bool): Whether sent to Telegram
    """
    import inspect
    import json

    topic = params.get('topic')
    if not topic:
        return {'success': False, 'error': 'topic parameter is required'}

    n_slides = params.get('n_slides', 10)
    template = params.get('template', 'dark')
    send_telegram = params.get('send_telegram', True)
    export_as = params.get('export_as', 'pptx').lower()  # 'pptx', 'pdf', or 'both'

    try:
        try:
            from Jotty.core.registry.skills_registry import get_skills_registry
        except ImportError:
            from Jotty.core.registry.skills_registry import get_skills_registry

        registry = get_skills_registry()
        registry.init()

        claude_skill = registry.get_skill('claude-cli-llm')
        telegram_skill = registry.get_skill('telegram-sender')

        if not claude_skill:
            return {'success': False, 'error': 'claude-cli-llm skill not available'}

        # Generate slide content with Claude
        prompt = f"""Create slide content for a presentation about: {topic}

Return ONLY a JSON object (no markdown, no explanation) with this exact structure:
{{
    "title": "Main Presentation Title",
    "subtitle": "Optional subtitle",
    "slides": [
        {{"title": "Slide 1 Title", "bullets": ["Point 1", "Point 2", "Point 3"]}},
        {{"title": "Slide 2 Title", "bullets": ["Point 1", "Point 2", "Point 3"]}}
    ]
}}

STRICT Requirements:
- Create exactly {n_slides} content slides
- Slide titles: Maximum 6 words, clear and concise
- Each slide: 3-5 bullet points
- CRITICAL: Each bullet MUST be 8-12 words maximum (shorter is better)
- Bullets should be punchy key points, not sentences
- No bullet should exceed 80 characters
- Make it informative and professional
- Cover the topic comprehensively

Return ONLY the JSON, nothing else."""

        generate_tool = claude_skill.tools.get('generate_text_tool')

        if inspect.iscoroutinefunction(generate_tool):
            result = await generate_tool({
                'prompt': prompt,
                'model': 'sonnet',
                'timeout': 120
            })
        else:
            result = generate_tool({
                'prompt': prompt,
                'model': 'sonnet',
                'timeout': 120
            })

        if not result.get('success'):
            return {'success': False, 'error': f"AI generation failed: {result.get('error')}"}

        # Parse JSON response
        text = result.get('text', '').strip()

        # Clean up response if needed
        if text.startswith('```'):
            text = text.split('```')[1]
            if text.startswith('json'):
                text = text[4:]
        text = text.strip()

        try:
            slide_data = json.loads(text)
        except json.JSONDecodeError as e:
            return {'success': False, 'error': f'Failed to parse AI response as JSON: {e}', 'raw_response': text[:500]}

        # Generate PPTX
        pptx_result = generate_slides_tool({
            'title': slide_data.get('title', topic),
            'subtitle': slide_data.get('subtitle', ''),
            'slides': slide_data.get('slides', []),
            'template': template,
            'output_path': params.get('output_path', os.path.expanduser('~/jotty/presentations'))
        })

        if not pptx_result.get('success'):
            return pptx_result

        pptx_path = pptx_result.get('file_path')
        pdf_path = None

        # Generate PDF directly if requested (using reportlab, not conversion)
        if export_as in ('pdf', 'both'):
            pdf_result = generate_slides_pdf_tool({
                'title': slide_data.get('title', topic),
                'subtitle': slide_data.get('subtitle', ''),
                'slides': slide_data.get('slides', []),
                'template': template,
                'output_path': params.get('output_path', os.path.expanduser('~/jotty/presentations'))
            })

            if pdf_result.get('success'):
                pdf_path = pdf_result.get('file_path')
                logger.info(f"PDF slides created: {pdf_path}")
            else:
                logger.warning(f"PDF generation failed: {pdf_result.get('error')}")

        # Determine which file to send to Telegram
        if export_as == 'pdf' and pdf_path:
            telegram_file = pdf_path
        else:
            telegram_file = pptx_path

        # Send to Telegram
        telegram_sent = False
        if send_telegram and telegram_skill:
            send_tool = telegram_skill.tools.get('send_telegram_file_tool')
            if send_tool:
                file_type = "PDF" if telegram_file.endswith('.pdf') else "PPTX"
                caption = f"📊 {slide_data.get('title', topic)}\n\n🎯 {n_slides + 1} slides\n📝 Topic: {topic}\n📄 Format: {file_type}"

                if inspect.iscoroutinefunction(send_tool):
                    tg_result = await send_tool({
                        'file_path': telegram_file,
                        'caption': caption
                    })
                else:
                    tg_result = send_tool({
                        'file_path': telegram_file,
                        'caption': caption
                    })

                telegram_sent = tg_result.get('success', False)

        result = {
            'success': True,
            'file_path': pptx_path,
            'slide_count': pptx_result.get('slide_count'),
            'title': slide_data.get('title'),
            'telegram_sent': telegram_sent,
            'export_as': export_as
        }

        if pdf_path:
            result['pdf_path'] = pdf_path

        return result

    except Exception as e:
        logger.error(f"Topic slide generation failed: {e}", exc_info=True)
        return {'success': False, 'error': f'Generation failed: {str(e)}'}
