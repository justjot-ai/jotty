"""
Theme Factory Skill - Apply professional themes to artifacts.

Provides curated color palettes and font pairings for consistent
styling across presentations, documents, and HTML pages.
"""
import asyncio
import logging
from pathlib import Path
from typing import Dict, Any, Optional
from datetime import datetime
import os

logger = logging.getLogger(__name__)

# Theme definitions
THEMES = {
    'ocean_depths': {
        'name': 'Ocean Depths',
        'description': 'Professional and calming maritime theme',
        'colors': {
            'primary': '#1a2332',
            'secondary': '#2d8b8b',
            'accent': '#a8dadc',
            'background': '#f1faee',
            'text': '#1a2332'
        },
        'fonts': {
            'heading': 'DejaVu Sans Bold',
            'body': 'DejaVu Sans'
        }
    },
    'sunset_boulevard': {
        'name': 'Sunset Boulevard',
        'description': 'Warm and vibrant sunset colors',
        'colors': {
            'primary': '#ff6b6b',
            'secondary': '#ffa07a',
            'accent': '#ffd93d',
            'background': '#fff5e6',
            'text': '#2c3e50'
        },
        'fonts': {
            'heading': 'Georgia Bold',
            'body': 'Georgia'
        }
    },
    'forest_canopy': {
        'name': 'Forest Canopy',
        'description': 'Natural and grounded earth tones',
        'colors': {
            'primary': '#2d5016',
            'secondary': '#5a7c3e',
            'accent': '#8fb98f',
            'background': '#f5f5dc',
            'text': '#2d5016'
        },
        'fonts': {
            'heading': 'Trebuchet MS Bold',
            'body': 'Trebuchet MS'
        }
    },
    'modern_minimalist': {
        'name': 'Modern Minimalist',
        'description': 'Clean and contemporary grayscale',
        'colors': {
            'primary': '#1a1a1a',
            'secondary': '#4a4a4a',
            'accent': '#808080',
            'background': '#ffffff',
            'text': '#1a1a1a'
        },
        'fonts': {
            'heading': 'Arial Bold',
            'body': 'Arial'
        }
    },
    'golden_hour': {
        'name': 'Golden Hour',
        'description': 'Rich and warm autumnal palette',
        'colors': {
            'primary': '#8b4513',
            'secondary': '#d2691e',
            'accent': '#f4a460',
            'background': '#fff8dc',
            'text': '#654321'
        },
        'fonts': {
            'heading': 'Times New Roman Bold',
            'body': 'Times New Roman'
        }
    },
    'arctic_frost': {
        'name': 'Arctic Frost',
        'description': 'Cool and crisp winter-inspired theme',
        'colors': {
            'primary': '#2c3e50',
            'secondary': '#5dade2',
            'accent': '#aed6f1',
            'background': '#ebf5fb',
            'text': '#1b2631'
        },
        'fonts': {
            'heading': 'Verdana Bold',
            'body': 'Verdana'
        }
    },
    'desert_rose': {
        'name': 'Desert Rose',
        'description': 'Soft and sophisticated dusty tones',
        'colors': {
            'primary': '#8b6f47',
            'secondary': '#c19a6b',
            'accent': '#e6d5b8',
            'background': '#faf8f3',
            'text': '#5d4e37'
        },
        'fonts': {
            'heading': 'Palatino Bold',
            'body': 'Palatino'
        }
    },
    'tech_innovation': {
        'name': 'Tech Innovation',
        'description': 'Bold and modern tech aesthetic',
        'colors': {
            'primary': '#0a0e27',
            'secondary': '#1e3a8a',
            'accent': '#3b82f6',
            'background': '#f8fafc',
            'text': '#0a0e27'
        },
        'fonts': {
            'heading': 'Courier New Bold',
            'body': 'Courier New'
        }
    },
    'botanical_garden': {
        'name': 'Botanical Garden',
        'description': 'Fresh and organic garden colors',
        'colors': {
            'primary': '#2d5016',
            'secondary': '#4a7c59',
            'accent': '#90c695',
            'background': '#f0f7f4',
            'text': '#1a3d1a'
        },
        'fonts': {
            'heading': 'Garamond Bold',
            'body': 'Garamond'
        }
    },
    'midnight_galaxy': {
        'name': 'Midnight Galaxy',
        'description': 'Dramatic and cosmic deep tones',
        'colors': {
            'primary': '#0d0221',
            'secondary': '#2d1b69',
            'accent': '#6b46c1',
            'background': '#f3e8ff',
            'text': '#0d0221'
        },
        'fonts': {
            'heading': 'Impact',
            'body': 'Arial'
        }
    }
}

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available, PPTX styling disabled")


async def apply_theme_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Apply a theme to an artifact.
    
    Args:
        params:
            - theme_name (str): Theme name
            - artifact_path (str): Path to artifact
            - artifact_type (str, optional): Type of artifact
            - custom_colors (dict, optional): Custom colors
            - custom_fonts (dict, optional): Custom fonts
    
    Returns:
        Dictionary with theme details and output path
    """
    theme_name = params.get('theme_name', '')
    artifact_path = params.get('artifact_path', '')
    artifact_type = params.get('artifact_type', 'auto')
    custom_colors = params.get('custom_colors', {})
    custom_fonts = params.get('custom_fonts', {})
    
    if not theme_name:
        return {
            'success': False,
            'error': 'theme_name is required'
        }
    
    if not artifact_path:
        return {
            'success': False,
            'error': 'artifact_path is required'
        }
    
    # Get theme
    if theme_name == 'custom':
        theme = {
            'name': 'Custom Theme',
            'description': 'User-defined custom theme',
            'colors': custom_colors or {},
            'fonts': custom_fonts or {}
        }
    elif theme_name in THEMES:
        theme = THEMES[theme_name]
    else:
        return {
            'success': False,
            'error': f'Unknown theme: {theme_name}. Available: {", ".join(THEMES.keys())}'
        }
    
    artifact_file = Path(os.path.expanduser(artifact_path))
    if not artifact_file.exists():
        return {
            'success': False,
            'error': f'Artifact file not found: {artifact_path}'
        }
    
    # Auto-detect artifact type
    if artifact_type == 'auto':
        suffix = artifact_file.suffix.lower()
        if suffix == '.pptx':
            artifact_type = 'pptx'
        elif suffix in ['.html', '.htm']:
            artifact_type = 'html'
        elif suffix == '.css':
            artifact_type = 'css'
        else:
            return {
                'success': False,
                'error': f'Unsupported artifact type: {suffix}. Supported: .pptx, .html, .css'
            }
    
    # Determine output path
    stem = artifact_file.stem
    suffix = artifact_file.suffix
    output_path = artifact_file.parent / f"{stem}_themed{suffix}"
    
    try:
        if artifact_type == 'pptx':
            result = await _apply_theme_to_pptx(artifact_file, output_path, theme)
        elif artifact_type == 'html':
            result = await _apply_theme_to_html(artifact_file, output_path, theme)
        elif artifact_type == 'css':
            result = await _apply_theme_to_css(artifact_file, output_path, theme)
        else:
            return {
                'success': False,
                'error': f'Unsupported artifact type: {artifact_type}'
            }
        
        return {
            'success': True,
            'theme_applied': theme['name'],
            'colors': theme['colors'],
            'fonts': theme['fonts'],
            'output_path': str(output_path)
        }
        
    except Exception as e:
        logger.error(f"Theme application failed: {e}", exc_info=True)
        return {
            'success': False,
            'error': str(e)
        }


def _hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


async def _apply_theme_to_pptx(input_path: Path, output_path: Path, theme: Dict) -> Dict:
    """Apply theme to PowerPoint presentation."""
    
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx not available. Install with: pip install python-pptx")
    
    prs = Presentation(str(input_path))
    
    for slide in prs.slides:
        for shape in slide.shapes:
            # Style text
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Apply font
                        font_size = run.font.size
                        if font_size and font_size.pt >= 24:
                            run.font.name = theme['fonts']['heading']
                        else:
                            run.font.name = theme['fonts']['body']
                        
                        # Apply text color
                        run.font.color.rgb = RGBColor(*_hex_to_rgb(theme['colors']['text']))
            
            # Style shapes
            if hasattr(shape, 'fill'):
                if hasattr(shape.fill, 'solid'):
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(theme['colors']['accent']))
    
    prs.save(str(output_path))
    return {}


async def _apply_theme_to_html(input_path: Path, output_path: Path, theme: Dict) -> Dict:
    """Apply theme to HTML file."""
    
    html_content = input_path.read_text(encoding='utf-8')
    
    css = f"""
    <style>
        body {{
            font-family: '{theme['fonts']['body']}', serif;
            color: {theme['colors']['text']};
            background-color: {theme['colors']['background']};
        }}
        h1, h2, h3, h4, h5, h6 {{
            font-family: '{theme['fonts']['heading']}', sans-serif;
            color: {theme['colors']['primary']};
        }}
        .accent {{
            color: {theme['colors']['accent']};
        }}
        .secondary {{
            color: {theme['colors']['secondary']};
        }}
    </style>
    """
    
    if '</head>' in html_content:
        html_content = html_content.replace('</head>', css + '</head>')
    else:
        html_content = css + html_content
    
    output_path.write_text(html_content, encoding='utf-8')
    return {}


async def _apply_theme_to_css(input_path: Path, output_path: Path, theme: Dict) -> Dict:
    """Apply theme to CSS file."""
    
    css_content = input_path.read_text(encoding='utf-8')
    
    theme_css = f"""
    /* Theme: {theme['name']} */
    :root {{
        --primary-color: {theme['colors']['primary']};
        --secondary-color: {theme['colors']['secondary']};
        --accent-color: {theme['colors']['accent']};
        --background-color: {theme['colors']['background']};
        --text-color: {theme['colors']['text']};
        --heading-font: '{theme['fonts']['heading']}';
        --body-font: '{theme['fonts']['body']}';
    }}
    
    body {{
        font-family: var(--body-font);
        color: var(--text-color);
        background-color: var(--background-color);
    }}
    
    h1, h2, h3, h4, h5, h6 {{
        font-family: var(--heading-font);
        color: var(--primary-color);
    }}
    
    {css_content}
    """
    
    output_path.write_text(theme_css, encoding='utf-8')
    return {}
