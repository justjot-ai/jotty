"""
Brand Guidelines Skill - Apply Anthropic brand styling.

Applies official brand colors and typography to documents
and presentations for consistent brand identity.
"""
import asyncio
import logging
from pathlib import Path
from typing import Dict, Any, Optional
from datetime import datetime
import os

from Jotty.core.infrastructure.utils.skill_status import SkillStatus
from Jotty.core.infrastructure.utils.tool_helpers import tool_response, tool_error, async_tool_wrapper

logger = logging.getLogger(__name__)

# Brand Colors
BRAND_COLORS = {
    'dark': '#141413',
    'light': '#faf9f5',
    'mid_gray': '#b0aea5',
    'light_gray': '#e8e6dc',
    'orange': '#d97757',
    'blue': '#6a9bcc',
    'green': '#788c5d'
}

# Typography

# Status emitter for progress updates
status = SkillStatus("brand-guidelines")

BRAND_FONTS = {
    'heading': 'Poppins',  # Fallback: Arial
    'body': 'Lora'  # Fallback: Georgia
}

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available, PPTX styling disabled")

try:
    from docx import Document
    from docx.shared import RGBColor as DocxRGBColor, Pt as DocxPt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available, DOCX styling disabled")


@async_tool_wrapper()
async def apply_brand_styling_tool(params: Dict[str, Any]) -> Dict[str, Any]:
    """
    Apply brand styling to a document or artifact.
    
    Args:
        params:
            - input_file (str): Path to input file
            - output_file (str, optional): Output path
            - file_type (str, optional): File type
            - preserve_content (bool, optional): Preserve content
    
    Returns:
        Dictionary with styled output path and applied styles
    """
    status.set_callback(params.pop('_status_callback', None))

    input_file = params.get('input_file', '')
    output_file = params.get('output_file', None)
    file_type = params.get('file_type', None)
    preserve_content = params.get('preserve_content', True)
    
    if not input_file:
        return {
            'success': False,
            'error': 'input_file is required'
        }
    
    input_path = Path(os.path.expanduser(input_file))
    if not input_path.exists():
        return {
            'success': False,
            'error': f'Input file not found: {input_file}'
        }
    
    # Auto-detect file type
    if not file_type:
        suffix = input_path.suffix.lower()
        if suffix == '.pptx':
            file_type = 'pptx'
        elif suffix == '.docx':
            file_type = 'docx'
        elif suffix in ['.html', '.htm']:
            file_type = 'html'
        else:
            return {
                'success': False,
                'error': f'Unsupported file type: {suffix}. Supported: .pptx, .docx, .html'
            }
    
    # Determine output path
    if not output_file:
        stem = input_path.stem
        suffix = input_path.suffix
        output_file = input_path.parent / f"{stem}_branded{suffix}"
    else:
        output_file = Path(os.path.expanduser(output_file))
    
    styles_applied = {}
    
    try:
        if file_type == 'pptx':
            if not PPTX_AVAILABLE:
                return {
                    'success': False,
                    'error': 'python-pptx not available. Install with: pip install python-pptx'
                }
            
            result = await _style_pptx(input_path, output_file)
            styles_applied = result.get('styles_applied', {})
        
        elif file_type == 'docx':
            if not DOCX_AVAILABLE:
                return {
                    'success': False,
                    'error': 'python-docx not available. Install with: pip install python-docx'
                }
            result = await _style_docx(input_path, output_file)
            styles_applied = result.get('styles_applied', {})
        
        elif file_type == 'html':
            result = await _style_html(input_path, output_file)
            styles_applied = result.get('styles_applied', {})
        
        else:
            return {
                'success': False,
                'error': f'Unsupported file type: {file_type}'
            }
        
        return {
            'success': True,
            'input_file': str(input_path),
            'output_file': str(output_file),
            'styles_applied': styles_applied
        }
        
    except Exception as e:
        logger.error(f"Brand styling failed: {e}", exc_info=True)
        return {
            'success': False,
            'error': str(e)
        }


def _hex_to_rgb(hex_color: str) -> tuple:
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))


async def _style_pptx(input_path: Path, output_path: Path) -> Dict:
    """Apply brand styling to PowerPoint presentation."""
    
    prs = Presentation(str(input_path))
    styles_applied = {
        'slides_styled': 0,
        'colors_applied': [],
        'fonts_applied': []
    }
    
    accent_colors = [BRAND_COLORS['orange'], BRAND_COLORS['blue'], BRAND_COLORS['green']]
    accent_index = 0
    
    for slide in prs.slides:
        styles_applied['slides_styled'] += 1
        
        for shape in slide.shapes:
            # Style text
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Apply font
                        font_size = run.font.size
                        if font_size and font_size.pt >= 24:
                            # Heading
                            run.font.name = BRAND_FONTS['heading']
                            styles_applied['fonts_applied'].append('Poppins (headings)')
                        else:
                            # Body
                            run.font.name = BRAND_FONTS['body']
                            styles_applied['fonts_applied'].append('Lora (body)')
                        
                        # Apply colors
                        run.font.color.rgb = RGBColor(*_hex_to_rgb(BRAND_COLORS['dark']))
            
            # Style shapes with accent colors
            if hasattr(shape, 'fill'):
                if hasattr(shape.fill, 'solid'):
                    hex_color = accent_colors[accent_index % len(accent_colors)]
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(hex_color))
                    styles_applied['colors_applied'].append(hex_color)
                    accent_index += 1
    
    prs.save(str(output_path))
    
    return {
        'styles_applied': styles_applied
    }


async def _style_html(input_path: Path, output_path: Path) -> Dict:
    """Apply brand styling to HTML file."""
    
    html_content = input_path.read_text(encoding='utf-8')
    
    # Inject CSS styles
    css = f"""
    <style>
        body {{
            font-family: '{BRAND_FONTS['body']}', Georgia, serif;
            color: {BRAND_COLORS['dark']};
            background-color: {BRAND_COLORS['light']};
        }}
        h1, h2, h3, h4, h5, h6 {{
            font-family: '{BRAND_FONTS['heading']}', Arial, sans-serif;
            color: {BRAND_COLORS['dark']};
        }}
        .accent-orange {{ color: {BRAND_COLORS['orange']}; }}
        .accent-blue {{ color: {BRAND_COLORS['blue']}; }}
        .accent-green {{ color: {BRAND_COLORS['green']}; }}
    </style>
    """
    
    # Insert CSS before </head> or at the beginning
    if '</head>' in html_content:
        html_content = html_content.replace('</head>', css + '</head>')
    else:
        html_content = css + html_content
    
    output_path.write_text(html_content, encoding='utf-8')

    return {
        'styles_applied': {
            'colors_applied': list(BRAND_COLORS.values()),
            'fonts_applied': list(BRAND_FONTS.values())
        }
    }


async def _style_docx(input_path: Path, output_path: Path) -> Dict:
    """Apply brand styling to Word document."""

    doc = Document(str(input_path))
    styles_applied = {
        'paragraphs_styled': 0,
        'colors_applied': [],
        'fonts_applied': []
    }

    # Style paragraphs
    for paragraph in doc.paragraphs:
        styles_applied['paragraphs_styled'] += 1

        for run in paragraph.runs:
            # Determine if heading or body based on style or font size
            is_heading = (
                paragraph.style.name.startswith('Heading') or
                (run.font.size and run.font.size.pt >= 14)
            )

            # Apply font
            if is_heading:
                run.font.name = BRAND_FONTS['heading']
                if 'Poppins (headings)' not in styles_applied['fonts_applied']:
                    styles_applied['fonts_applied'].append('Poppins (headings)')
            else:
                run.font.name = BRAND_FONTS['body']
                if 'Lora (body)' not in styles_applied['fonts_applied']:
                    styles_applied['fonts_applied'].append('Lora (body)')

            # Apply text color
            run.font.color.rgb = DocxRGBColor(*_hex_to_rgb(BRAND_COLORS['dark']))

    # Style tables if present
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = BRAND_FONTS['body']
                        run.font.color.rgb = DocxRGBColor(*_hex_to_rgb(BRAND_COLORS['dark']))

    styles_applied['colors_applied'] = [BRAND_COLORS['dark']]

    doc.save(str(output_path))

    return {
        'styles_applied': styles_applied
    }
